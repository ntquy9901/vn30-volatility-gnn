"""
make_slides.py — VN30 Volatility Forecasting — Báo cáo luận văn
Slides: Title | Overview | Arch×2 | Results×2 | Findings×5 | Pretraining | Conclusions | Next
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x37, 0x6C)
ACCENT_BLUE = RGBColor(0x27, 0x6E, 0xB4)
ACCENT_GOLD = RGBColor(0xE6, 0x9B, 0x1A)
LIGHT_GRAY  = RGBColor(0xF2, 0xF4, 0xF8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
GREEN       = RGBColor(0x1A, 0x7A, 0x3C)
GREEN_DK    = RGBColor(0x0D, 0x5C, 0x2A)
RED_DARK    = RGBColor(0xAA, 0x22, 0x22)
PURPLE      = RGBColor(0x5B, 0x0E, 0x91)
PURPLE_DK   = RGBColor(0x3E, 0x00, 0x6B)
ORANGE      = RGBColor(0xC4, 0x52, 0x00)
ORANGE_DK   = RGBColor(0x93, 0x3C, 0x00)
BG          = RGBColor(0xF8, 0xF9, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Primitive helpers ─────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0.75)):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh


def txt(slide, text, l, t, w, h,
        size=12, bold=False, color=DARK_GRAY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def bullets(slide, items, l, t, w, h, size=12, color=DARK_GRAY):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(2)
        r = p.add_run(); r.text = item
        r.font.size = Pt(size); r.font.color.rgb = color


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, fill=DARK_BLUE)
    txt(slide, title, 0.3, 0.07, 12.6, 0.72,
        size=25, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.75, 12.6, 0.35,
            size=12, color=RGBColor(0xBB, 0xCC, 0xEE), align=PP_ALIGN.LEFT)


def table(slide, headers_row, rows, l, t, w,
          col_widths=None, row_h=0.38,
          hdr_fill=DARK_BLUE, alt=LIGHT_GRAY,
          hdr_size=10, data_size=9):
    n = len(headers_row)
    if col_widths is None:
        col_widths = [w / n] * n
    x = l
    for hd, cw in zip(headers_row, col_widths):
        rect(slide, x, t, cw, row_h, fill=hdr_fill)
        txt(slide, hd, x+0.04, t+0.04, cw-0.08, row_h-0.08,
            size=hdr_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw
    for ri, row in enumerate(rows):
        y = t + row_h * (ri + 1)
        x = l
        fill = alt if ri % 2 == 1 else WHITE
        for cell, cw in zip(row, col_widths):
            rect(slide, x, y, cw, row_h, fill=fill,
                 line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.4))
            txt(slide, str(cell), x+0.03, y+0.03, cw-0.06, row_h-0.06,
                size=data_size, color=DARK_GRAY, align=PP_ALIGN.CENTER)
            x += cw


def img(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


def arrow_v(slide, x, y, color=DARK_GRAY):
    """Small downward arrow symbol."""
    txt(slide, "▼", x-0.15, y, 0.3, 0.32, size=13, color=color, align=PP_ALIGN.CENTER)


def arch_box(slide, l, t, w, h, title, body_lines, fill, title_size=12, body_size=10.5):
    rect(slide, l, t, w, h, fill=fill)
    txt(slide, title, l+0.08, t+0.07, w-0.16, 0.38,
        size=title_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(slide, body_lines, l+0.08, t+0.48, w-0.16, h-0.55,
            size=body_size, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
rect(s, 0, 5.8, 13.33, 1.7, fill=ACCENT_BLUE)
txt(s, "DU BAO BIEN DONG THUC HIEN (REALIZED VOLATILITY)",
    0.5, 1.2, 12.3, 0.9, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Co phieu VN30 — Danh gia Salesforce Moirai2 & Graph Neural Network",
    0.5, 2.1, 12.3, 0.65, size=19, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
rect(s, 4.5, 3.0, 4.33, 0.06, fill=ACCENT_GOLD)
txt(s, "Du lieu: VN30, 2019-2026  |  Test set: 2026-01-05 -> 2026-04-07 (62 ngay, 30 co phieu, 1,860 cap)",
    0.5, 3.3, 12.3, 0.5, size=13, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
txt(s, "Luan van tot nghiep  —  2026-05-16",
    0.5, 6.0, 12.3, 0.5, size=13, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "HCMUS  —  Khoa Cong nghe Thong tin",
    0.5, 6.5, 12.3, 0.5, size=12, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Tổng quan
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Tong Quan Nghien Cuu", "Muc tieu, du lieu, cau hoi nghien cuu")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

rect(s, 0.3, 1.28, 5.8, 2.75, fill=WHITE, line=ACCENT_BLUE, line_w=Pt(1.5))
txt(s, "Muc Tieu", 0.5, 1.33, 5.4, 0.38, size=13, bold=True, color=ACCENT_BLUE)
bullets(s, [
    "Danh gia Moirai2 (foundation model) cho du bao RV tren VN30",
    "So sanh: embedding vs IQR quantile spread vs HAR-RV (classical)",
    "Can thiep gi cua foundation model mang lai gia tri cho bai toan nay?",
    "Tim cach thu hep khoang cach giua neural models va HAR-RV",
], 0.5, 1.73, 5.6, 2.2, size=12.5)

rect(s, 6.5, 1.28, 6.5, 2.75, fill=WHITE, line=ACCENT_BLUE, line_w=Pt(1.5))
txt(s, "Du Lieu & Thiet Ke", 6.7, 1.33, 6.2, 0.38, size=13, bold=True, color=ACCENT_BLUE)
bullets(s, [
    "30 co phieu VN30, OHLCV daily, ~2,500 ngay (2019-2026)",
    "Train: full history per-stock HOAC batch 2,896 mau (tat ca stocks)",
    "Test: 62 ngay giao dich, 01/2026-04/2026 = 1,860 cap (date, stock)",
    "Nhan (label): RV h=20 = std(log-return, next 20 ngay)",
    "Metrics danh gia: MAE, R2, Pearson_r, QLIKE",
], 6.7, 1.73, 6.2, 2.2, size=12.5)

rect(s, 0.3, 4.15, 12.7, 3.0, fill=WHITE, line=ACCENT_BLUE, line_w=Pt(1.5))
txt(s, "12 Mo Hinh Duoc Danh Gia (3 nhom chinh)", 0.5, 4.2, 12.3, 0.38, size=13, bold=True, color=ACCENT_BLUE)
bullets(s, [
    "Classical baselines: HAR-RV (OLS per-stock), LSTM (per-stock), GARCH(1,1) (MLE per-stock)",
    "Moirai2 zero-shot: (A) Embedding 384-dim → MLP/GNN   (B) IQR = q(0.9) - q(0.1) lam proxy RV",
    "Neural (walk-forward): MLP+Moirai2, GNN+Moirai2 [ham chuan du lieu 30 mau/buoc]",
    "Neural (batch): MLP(RV6), GNN(RV6), MLP(Moirai2+RV6), GNN(Moirai2+RV6), MLP/GNN(Moirai2+RV3)",
], 0.5, 4.6, 12.5, 2.4, size=12.5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture 1: Overall Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Kien Truc Thi Nghiem — Tong The Pipeline",
       "Ba nhanh song song: Classical | Moirai2+Neural | GARCH")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

# ── Row 1: Data input ─────────────────────────────────────────────────────────
rect(s, 0.3, 1.22, 12.7, 0.6, fill=DARK_BLUE)
txt(s, "VN30 OHLCV Data  —  30 co phieu  |  ~2,500 ngay giao dich (2019-2026)  |  Close price daily",
    0.5, 1.33, 12.3, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Arrow + preprocessing ─────────────────────────────────────────────────────
arrow_v(s, 6.66, 1.84)
rect(s, 4.1, 2.17, 5.1, 0.52, fill=ACCENT_BLUE)
txt(s, "Log Returns:  r_t = log( Close_t / Close_{t-1} )",
    4.22, 2.26, 4.86, 0.34, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Branch arrows ─────────────────────────────────────────────────────────────
txt(s, "↙", 3.6, 2.72, 0.5, 0.38, size=20, color=GREEN, align=PP_ALIGN.CENTER)
txt(s, "↓", 6.55, 2.72, 0.4, 0.38, size=20, color=PURPLE, align=PP_ALIGN.CENTER)
txt(s, "↘", 9.2, 2.72, 0.5, 0.38, size=20, color=ORANGE, align=PP_ALIGN.CENTER)

# ── Column layout ─────────────────────────────────────────────────────────────
C1X, C1W = 0.25, 3.85   # Classical
C2X, C2W = 4.55, 4.45   # Moirai2
C3X, C3W = 9.45, 3.6    # GARCH

# ── Row 2: Feature extraction ─────────────────────────────────────────────────
arch_box(s, C1X, 3.12, C1W, 1.25, "RV Features (6-dim)", [
    "log(RV_d)  log(RV_w)  log(RV_m)",
    "log(RV_q=60d)  corr_vnindex",
    "jump_ratio",
    "(Corsi 2009 + mo rong)",
], fill=GREEN, body_size=10.5)

arch_box(s, C2X, 3.12, C2W, 1.25, "Moirai2-Small  (FROZEN, 11.4M params)", [
    "Context 200 ngay  →  14 patches  →  Decoder Transformer",
    "Output A: Embedding 384-dim (patch cuoi cung, pos -2)",
    "Output B: IQR = q(0.9) - q(0.1)  [zero-shot proxy RV]",
], fill=PURPLE, body_size=10.5)

arch_box(s, C3X, 3.12, C3W, 1.25, "GARCH(1,1)", [
    "MLE per-stock",
    "s^2_t = w + a*e^2_{t-1}",
    "         + b*s^2_{t-1}",
    "(benchmark truyen thong)",
], fill=ORANGE, body_size=10.5)

# ── Arrows down ───────────────────────────────────────────────────────────────
arrow_v(s, C1X + C1W/2, 4.4, color=GREEN)
arrow_v(s, C2X + C2W/2, 4.4, color=PURPLE)
arrow_v(s, C3X + C3W/2, 4.4, color=ORANGE)

# ── Row 3: Models ─────────────────────────────────────────────────────────────
rect(s, C1X, 4.73, C1W, 1.35, fill=GREEN_DK)
txt(s, "HAR-RV (OLS)  /  LSTM", C1X+0.08, 4.78, C1W-0.16, 0.38,
    size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "BEST", C1X+C1W-1.1, 4.76, 0.95, 0.38,
    size=12, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.RIGHT)
bullets(s, [
    "Per-stock, full history (~2,500 mau/stock)",
    "HAR: 3 tham so OLS (BLUE estimator)",
    "MAE: 0.0010 / 0.0011   R2: 0.971 / 0.967",
], C1X+0.08, 5.18, C1W-0.16, 0.85, size=10.5, color=WHITE)

rect(s, C2X, 4.73, C2W, 1.35, fill=PURPLE_DK)
txt(s, "MLP / GNN  (batch or walk-forward)", C2X+0.08, 4.78, C2W-0.16, 0.38,
    size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s, [
    "Input: Emb(384)+RV3/RV6 | IQR zero-shot | RV6 only",
    "Batch 2,896 mau | Walk-fwd 30 mau/buoc",
    "MAE: 0.0066-0.0143   R2: -1.86 to +0.28",
], C2X+0.08, 5.18, C2W-0.16, 0.85, size=10.5, color=WHITE)

rect(s, C3X, 4.73, C3W, 1.35, fill=ORANGE_DK)
txt(s, "GARCH(1,1)", C3X+0.08, 4.78, C3W-0.16, 0.38,
    size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(s, [
    "MLE per-stock",
    "MAE: 0.0071   R2: +0.159",
], C3X+0.08, 5.18, C3W-0.16, 0.85, size=10.5, color=WHITE)

# ── Output ────────────────────────────────────────────────────────────────────
arrow_v(s, 6.66, 6.1, color=DARK_GRAY)
rect(s, 3.0, 6.43, 7.3, 0.55, fill=ACCENT_GOLD)
txt(s, "Realized Volatility Prediction  |  h = 20 ngay  =  std(log-return, next 20 ngay)",
    3.12, 6.52, 7.06, 0.38, size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture 2: Moirai2 Internals
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Kien Truc Moirai2 — Hai Che Do Su Dung",
       "Frozen decoder-only transformer: (A) Embedding extraction  vs  (B) IQR zero-shot")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

# ── Input series ──────────────────────────────────────────────────────────────
rect(s, 0.3, 1.25, 12.7, 0.55, fill=ACCENT_BLUE)
txt(s, "Log Return Series  —  200 observations (context window)  —  daily frequency",
    0.5, 1.35, 12.3, 0.35, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Patcher ───────────────────────────────────────────────────────────────────
arrow_v(s, 6.66, 1.82)
rect(s, 2.5, 2.16, 8.3, 0.52, fill=RGBColor(0x34, 0x5C, 0x8A))
txt(s, "Patcher:  200 obs / patch_size=16  →  13 context patches  +  1 MASK token  =  14 patches",
    2.65, 2.25, 8.0, 0.34, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Patch visualization ───────────────────────────────────────────────────────
arrow_v(s, 6.66, 2.7)
PATCH_Y = 3.05; PATCH_W = 0.6; PATCH_H = 0.5; GAP = 0.65
START_X = 0.55
for i in range(13):
    fill = ACCENT_BLUE if i < 12 else GREEN
    rect(s, START_X + i * GAP, PATCH_Y, PATCH_W, PATCH_H, fill=fill)
    txt(s, f"P{i+1}", START_X + i * GAP + 0.05, PATCH_Y + 0.08, 0.5, 0.32,
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# MASK token
MASK_X = START_X + 13 * GAP
rect(s, MASK_X, PATCH_Y, PATCH_W, PATCH_H, fill=RGBColor(0x88, 0x88, 0x88))
txt(s, "MASK", MASK_X + 0.03, PATCH_Y + 0.08, 0.54, 0.32,
    size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Labels
txt(s, "13 context patches (causal attention)", 0.5, PATCH_Y + 0.52, 8.5, 0.32,
    size=10, italic=True, color=RGBColor(0x44, 0x66, 0xAA))
txt(s, "← predict", MASK_X - 0.05, PATCH_Y + 0.52, 1.4, 0.32,
    size=10, italic=True, color=RGBColor(0x66, 0x66, 0x66))

# ── Transformer block ─────────────────────────────────────────────────────────
arrow_v(s, 6.66, 3.6)
rect(s, 2.0, 3.93, 9.3, 0.55, fill=DARK_BLUE)
txt(s, "Decoder-only Transformer  (causal masked attention, 11.4M params)  —  FROZEN trong toan bo thi nghiem",
    2.15, 4.02, 9.0, 0.35, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Split into 2 approaches ───────────────────────────────────────────────────
txt(s, "↙", 4.8, 4.5, 0.5, 0.38, size=20, color=GREEN, align=PP_ALIGN.CENTER)
txt(s, "↘", 8.0, 4.5, 0.5, 0.38, size=20, color=PURPLE, align=PP_ALIGN.CENTER)

# Approach A: Embedding
rect(s, 0.3, 4.9, 6.0, 1.9, fill=GREEN_DK)
txt(s, "Approach A — Embedding Extraction", 0.45, 4.95, 5.7, 0.4,
    size=13, bold=True, color=WHITE)
bullets(s, [
    "Lay hidden state tai patch cuoi cung (pos -2: last_context)",
    "Kich thuoc: 384-dim per stock per window",
    "→ Dung lam node features cho MLP / GNN",
    "Ket qua: median|corr(emb, RV)| = 0.128",
    "  Them vao RV6: MAE tang 8% (them nhieu!)",
], 0.45, 5.38, 5.7, 1.35, size=11, color=WHITE)

# Approach B: IQR
rect(s, 7.0, 4.9, 6.0, 1.9, fill=PURPLE_DK)
txt(s, "Approach B — IQR Quantile Spread (Zero-shot)", 7.15, 4.95, 5.7, 0.4,
    size=13, bold=True, color=WHITE)
bullets(s, [
    "Lay output du bao: QuantileForecast (9 quantile levels)",
    "q0.1, q0.2, q0.3, ..., q0.9  (exact, khong sampling)",
    "IQR_t = q(0.9) - q(0.1)  [80% prediction interval width]",
    "Scale: IQR x 0.442 ≈ RV (mean(RV)/mean(IQR))",
    "  Pearson_r = 0.557 — bang MLP trained (r=0.551)!",
], 7.15, 5.38, 5.7, 1.35, size=11, color=WHITE)

# Verdict labels
rect(s, 0.3, 6.82, 6.0, 0.45, fill=RED_DARK)
txt(s, "Embedding weak — domain gap (no equity pretraining)",
    0.45, 6.89, 5.7, 0.3, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 7.0, 6.82, 6.0, 0.45, fill=GREEN)
txt(s, "IQR competitive — uncertainty calibration hoat dong!",
    7.15, 6.89, 5.7, 0.3, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Full Results Table (all 12 models)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Ket Qua — Tat Ca 12 Mo Hinh", "Test set: 62 ngay x 30 co phieu = 1,860 cap | sap xep theo MAE tang dan")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

H = ["Mo hinh", "Training", "Features", "MAE", "R2", "Pearson_r", "QLIKE"]
CW = [3.2, 1.75, 1.75, 1.05, 1.05, 1.35, 1.05]
ROWS12 = [
    ["HAR-RV (OLS) ***",         "Full/stock",  "RV_d,w,m",   "0.0010", "+0.971", "0.985", "0.012"],
    ["LSTM ***",                  "Full/stock",  "same",       "0.0011", "+0.967", "0.983", "0.012"],
    ["MLP (RV6 only) **",        "Batch 2896",  "6 RV ext.",  "0.0066", "+0.278", "0.551", "0.328"],
    ["Moirai2 IQR (zero-shot)",  "Khong",       "q0.9-q0.1",  "0.0067+","+0.161", "0.557", "  — "],
    ["GARCH(1,1)",               "Per-stk MLE", "—",          "0.0071", "+0.159", "0.535", "0.301"],
    ["MLP (Moirai2+RV6)",        "Batch",       "390-dim",    "0.0072", "+0.065", "0.505", "0.259"],
    ["Batch GNN (Moirai2+RV3)", "Batch",        "387-dim",    "0.0074", "+0.041", "0.357", "0.399"],
    ["Batch MLP (Moirai2+RV3)", "Batch",        "387-dim",    "0.0078", "-0.137", "0.425", "0.290"],
    ["GNN (RV6 only) **",       "Batch",        "6 RV ext.",  "0.0084", "-0.194", "0.396", "0.328"],
    ["GNN (Moirai2+RV6)",       "Batch",        "390-dim",    "0.0091", "-0.519", "0.377", "0.345"],
    ["WalkFwd MLP+Moirai2",     "Walk-fwd",     "387-dim",    "0.0121", "-1.077", "0.239", "1.444"],
    ["WalkFwd GNN+Moirai2",     "Walk-fwd",     "387-dim",    "0.0143", "-1.856", "0.117", "1.898"],
]
table(s, H, ROWS12, 0.25, 1.25, 11.2, col_widths=CW, row_h=0.395,
      hdr_size=11, data_size=9.5)

# Legend
rect(s, 0.25, 6.22, 11.2, 0.95, fill=WHITE, line=ACCENT_BLUE, line_w=Pt(0.75))
bullets(s, [
    "*** Per-stock full history training (~2,500 mau/stock, OLS 3 params)   ** Batch 2,896 mau chung 30 stocks",
    "+ MAE sau scale correction (IQR x 0.442); raw MAE = 0.0323   | Bold = best per metric",
    "Walk-fwd = 30 mau/buoc (nut that training)  |  Batch = 2,896 mau tong",
], 0.4, 6.28, 11.0, 0.82, size=10)

# Metric note on right
rect(s, 11.6, 1.25, 1.5, 5.92, fill=RGBColor(0xF0, 0xF4, 0xFF), line=ACCENT_BLUE, line_w=Pt(0.75))
txt(s, "Best:", 11.7, 1.3, 1.3, 0.3, size=10, bold=True, color=DARK_BLUE)
bullets(s, [
    "MAE↓",  "0.0010",
    "R2↑",   "+0.971",
    "r↑",    "0.985",
    "QLIKE↓","0.012",
    "",
    "HAR-RV",
    "wins all",
], 11.7, 1.62, 1.3, 5.4, size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Results Visual (bar + R2 heatmap)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Ket Qua — Truc Quan Hoa",
       "Trai: MAE/RMSE/QLIKE so sanh 5 mo hinh chinh  |  Phai: R2 per stock x model")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)
img(s, "results/plot_bar_metrics.png", 0.2, 1.2, 7.0, 5.4)
img(s, "results/plot_r2_heatmap.png",  7.5, 1.2, 5.6, 6.1)
txt(s, "Chi HAR-RV va LSTM co R2 > 0 tren tat ca 30 co phieu (xanh la). Walk-fwd models: R2 am (do).",
    7.5, 7.2, 5.6, 0.35, size=9.5, italic=True,
    color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Finding 1: Training regime
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Finding 1 — Che Do Training La Nut That Chinh (H2)",
       "Walk-forward vs Batch: chenh lech 35-48% MAE")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

rect(s, 0.3, 1.28, 5.8, 4.6, fill=WHITE, line=RED_DARK, line_w=Pt(2))
txt(s, "Walk-Forward (che do cu)", 0.5, 1.33, 5.5, 0.4, size=14, bold=True, color=RED_DARK)
bullets(s, [
    "30 samples / gradient step",
    "(30 stocks x 1 ngay giao dich)",
    "",
    "MLP: MAE = 0.0121",
    "GNN: MAE = 0.0143",
    "",
    "Ly do: data starvation",
    "Moi window chi them 1 ngay",
    "Model khong the hoc duoc patterns",
], 0.5, 1.75, 5.5, 4.0, size=13)

txt(s, "vs", 6.2, 3.1, 0.8, 0.8, size=28, bold=True,
    color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

rect(s, 7.2, 1.28, 5.8, 4.6, fill=WHITE, line=GREEN, line_w=Pt(2))
txt(s, "Batch Training (cai thien)", 7.4, 1.33, 5.5, 0.4, size=14, bold=True, color=GREEN)
bullets(s, [
    "2,896 samples toan bo (tat ca stocks)",
    "(30 stocks, chung weights)",
    "",
    "MLP: MAE = 0.0066  (-45%)",
    "GNN: MAE = 0.0084  (-41%)",
    "",
    "Gop toan bo lich su 30 stocks",
    "Hoc duoc patterns chung",
    "Nhung van it hon HAR ~30x",
], 7.4, 1.75, 5.5, 4.0, size=13)

rect(s, 0.3, 5.98, 12.7, 1.18, fill=RGBColor(0xFF, 0xF3, 0xCD))
bullets(s, [
    "HAR-OLS: ~75,000 cap hoc (2,500 ngay x 30 stocks, fit rieng tung stock)  vs  Neural batch: ~97 mau/stock (2,896/30 shared)",
    "=> Gap 6.5x MAE la co cau truc — khong the giai quyet chi bang architecture phuc tap hon",
], 0.5, 6.03, 12.3, 1.05, size=12, color=RGBColor(0x66, 0x44, 0x00))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Finding 2 & 3: Embeddings + pooling chart
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Finding 2 & 3 — Embedding Them Nhieu | RV6 > RV3",
       "Domain gap giai thich tai sao median|corr(embedding,RV)|=0.128; RV6 giam 15% MAE")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

rect(s, 0.2, 1.25, 6.15, 2.05, fill=WHITE, line=RED_DARK, line_w=Pt(1.5))
txt(s, "Finding 2 — Embedding Moirai2 Them Nhieu (+8-9%)", 0.38, 1.3, 5.9, 0.38,
    size=12, bold=True, color=RED_DARK)
table(s, ["Mo hinh", "Features", "MAE"], [
    ["MLP (RV6 only)",    "6-dim RV",  "0.0066"],
    ["MLP (Moirai2+RV6)", "390-dim",   "0.0072 (+8.4%)"],
    ["GNN (RV6 only)",    "6-dim RV",  "0.0084"],
    ["GNN (Moirai2+RV6)", "390-dim",   "0.0091 (+8.1%)"],
], 0.3, 1.68, 5.95, col_widths=[2.1, 1.7, 2.0], row_h=0.39,
    hdr_fill=RED_DARK, data_size=10)

rect(s, 0.2, 3.38, 6.15, 3.27, fill=WHITE, line=GREEN, line_w=Pt(1.5))
txt(s, "Finding 3 — RV6 (6-dim) tot hon RV3 (3-dim): -15% MAE", 0.38, 3.43, 5.9, 0.38,
    size=12, bold=True, color=GREEN)
table(s, ["Feature", "Cong thuc", "Nguon"], [
    ["log(RV_d)",    "std(r, 1 ngay)",        "Corsi 2009"],
    ["log(RV_w)",    "mean(RVd, 5 ngay)",     "Corsi 2009"],
    ["log(RV_m)",    "mean(RVd, 20 ngay)",    "Corsi 2009"],
    ["log(RV_q)*",   "mean(RVd, 60 ngay)",   "Mo rong HAR"],
    ["corr_vnindex*","Pearson(stk,VNI,60d)", "Risk he thong"],
    ["jump_ratio*",  "max(RVd-RVw,0)/RVd",   "B-N&S 2004"],
], 0.3, 3.82, 5.95, col_widths=[1.75, 2.3, 1.75], row_h=0.38,
    hdr_fill=GREEN, data_size=10)
txt(s, "* 3 features mo rong ngoai HAR goc (RV_d, RV_w, RV_m)",
    0.38, 6.2, 5.7, 0.3, size=9.5, italic=True, color=RGBColor(0x55,0x55,0x55))

img(s, "results/diag_pooling_compare.png", 6.55, 1.22, 6.6, 4.55)
txt(s, "4 chien luoc pooling Moirai2 — tat ca deu yeu (median|corr| max 0.128)",
    6.55, 5.82, 6.6, 0.38, size=10.5, italic=True,
    color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)
rect(s, 6.55, 6.22, 6.6, 0.4, fill=RED_DARK)
txt(s, "Nguyen nhan co cau: Moirai2 chua bao gio thay equity/RV data trong pretraining",
    6.65, 6.28, 6.4, 0.28, size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Finding 5: IQR zero-shot + scatter
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Finding 5 — Moirai2 IQR: Zero-Shot Competitive",
       "IQR = q(0.9)-q(0.1) bat duoc volatility clustering — khong can training, khong can features")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

rect(s, 0.2, 1.22, 6.2, 0.62, fill=RGBColor(0xE8, 0xF4, 0xEC), line=GREEN, line_w=Pt(1.5))
txt(s, "IQR_t  =  q(0.9) - q(0.1)  of 1-step-ahead return forecast  →  proxy for RV",
    0.35, 1.32, 5.9, 0.42, size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

table(s, ["Phuong phap", "Pearson_r", "MAE", "Training can"],
[
    ["Moirai2 IQR (zero-shot)", "0.557",  "0.0067*", "Khong co"],
    ["MLP (RV6, batch)",        "0.551",  "0.0066",  "2,896 mau"],
    ["Embedding median|corr|",  "0.128",  "—",       "—"],
    ["HAR-RV (OLS)",            "0.985",  "0.0010",  "~2,500/stock"],
], 0.2, 1.9, 6.2, col_widths=[2.8, 1.3, 1.3, 1.7], row_h=0.44,
    hdr_fill=DARK_BLUE, data_size=10.5)

bullets(s, [
    "IQR khong can training, khong can feature engineering",
    "Pearson_r IQR (0.557) = 4.4x cao hon embedding (0.128)",
    "Scale factor 0.442 ≈ mean(RV)/mean(IQR)",
    "  => IQR o log-return units, lon hon RV ~2.26x",
    "Model hoc duoc heteroscedasticity (vol clustering)",
    "  du pretraining khong co equity/RV data",
    "* raw MAE = 0.0323; sau scale x0.442 = 0.0067",
], 0.2, 3.84, 6.2, 2.7, size=12)

img(s, "results/iqr_vs_rv.png", 6.55, 1.22, 6.6, 5.2)
txt(s, "Trai: IQR raw vs RV (r=0.557, R2=-13.4)  |  Phai: IQR scaled vs RV (MAE=0.0067)",
    6.55, 6.47, 6.6, 0.4, size=10, italic=True,
    color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Finding 6: HAR gap + timeseries
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Finding 6 — Khoang Cach 6.5x voi HAR-RV La Co Cau Truc",
       "HAR-RV theo sat true RV; neural models (walk-fwd) phang; batch models nhieu")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

rect(s, 0.2, 1.25, 5.6, 6.0, fill=WHITE, line=ACCENT_BLUE, line_w=Pt(1.5))
txt(s, "Tai sao HAR-RV vuot troi?", 0.4, 1.3, 5.3, 0.4, size=14, bold=True, color=ACCENT_BLUE)
bullets(s, [
    "Per-stock OLS: 2,500 mau x 30 stocks (fit rieng)",
    "Chi 3 tham so: [beta_d, beta_w, beta_m]",
    "MAE=0.0010  R2=+0.971  Pearson_r=0.985",
    "",
    "OLS = BLUE Estimator (Gauss-Markov theorem)",
    "Khong the co linear model nao tot hon",
    "voi cung du lieu",
    "",
    "Neural batch: ~97 mau/stock (2,896/30)",
    "=> MLP(RV6): MAE=0.0066 [gap 6.5x]",
    "",
    "Giai phap (lau dai):",
    "  HAR-RS: them asymmetry, van OLS",
    "  GNN+HAR: exploit cross-stock spillover",
    "  Mo rong data: them stocks, them nam",
], 0.4, 1.72, 5.3, 5.2, size=12)

img(s, "results/plot_timeseries.png", 6.0, 1.22, 7.1, 5.55)
txt(s, "VCB, HPG, FPT: HAR-RV (xanh la) bam sat true RV (den). Walk-fwd models: phang; LSTM: ok.",
    6.0, 6.82, 7.1, 0.4, size=10, italic=True,
    color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Pretraining Data
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Tai Sao Moirai2 Embedding Yeu? — Phan Tich Pretraining Data",
       "GIFT-Eval Pretrain corpus: 4.5M series, 230B diem — khong co equity data thuc su")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)
table(s, ["Dataset", "Tan suat", "So series", "Noi dung thuc su (giai thich chi tiet)"],
[
    ["M4 Daily",   "Daily",   "4,227",   "Business/macro competition — mot so stock, khong co Asian equity"],
    ["M4 Monthly", "Monthly", "48,000",  "Kinh te vi mo, demographic, industry sectors"],
    ["M1/M3",      "Mix",     "~3,000",  "Macro + business forecasting competition tong hop"],
    ["FRED-MD",    "Monthly", "107",     "Federal Reserve macro indicators (CPI, lai suat, employment) — chi My"],
    ["Bitcoin",    "Daily",   "18",      "Crypto price — DUY NHAT daily price series (khong phai stock equity)"],
    ["NN5",        "Daily",   "111",     "ATM cash withdrawal UK — bi label nham 'Econ/Fin', KHONG phai stock"],
    ["GoDaddy",    "Monthly", "3,135",   "Domain registration business data — khong lien quan den stock"],
], 0.3, 1.28, 12.7, col_widths=[1.8, 1.3, 1.5, 7.9], row_h=0.48,
    hdr_size=11, data_size=10)

rect(s, 0.3, 4.72, 12.7, 2.52, fill=RGBColor(0xFF, 0xEC, 0xEC), line=RED_DARK, line_w=Pt(1.5))
txt(s, "Ket Luan — Domain Gap La Nguyen Nhan Co Cau:", 0.5, 4.77, 12.3, 0.38,
    size=13, bold=True, color=RED_DARK)
bullets(s, [
    "KHONG co equity/stock market data thuc su — chi 18 series Bitcoin daily trong 4.5M series",
    "KHONG co Asian market data — toan bo corpus la US/EU-centric (FRED-MD, M4 competition)",
    "KHONG co realized volatility series — Moirai2 chua bao gio thay RV labels trong pretraining",
    "=> Embedding 384-dim phan anh macro trend/seasonality patterns, KHONG phai second-moment structure cua stock vol",
    "=> Day la giai thich co cau cho median|corr(embedding, RV)| = 0.128 tren tat ca 384 chieu",
], 0.5, 5.15, 12.3, 2.0, size=11.5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Kết luận
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Ket Luan", "5 insights chinh tu thi nghiem VN30 Volatility Forecasting")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

conclusions = [
    ("(1) Domain gap la nguyen nhan co cau",
     "Moirai2 pretrain tren macro/business data (M4, FRED-MD). Khong co equity/Asian/RV data. Embedding encode macro trend, khong phai second-moment structure => median|corr|=0.128."),
    ("(2) Training regime la nut that chinh",
     "Walk-forward (30 mau/buoc) → Batch (2,896 mau): giam 35-48% MAE. Training regime quan trong hon kien truc model."),
    ("(3) HAR-RV (OLS) van la benchmark champion",
     "R2=0.971 vs best neural R2=0.278. Gap 6.5x MAE la co cau: OLS per-stock 2,500+ mau, 3 params, BLUE estimator."),
    ("(4) IQR zero-shot — phat hien dac biet cua luan van",
     "Moirai2 IQR (q0.9-q0.1) dat Pearson_r=0.557 = MLP trained (r=0.551). Uncertainty calibration hoat dong du domain gap — model hoc duoc heteroscedasticity."),
    ("(5) Embedding Moirai2 them nhieu khi co RV features",
     "MLP(Moirai2+RV6) MAE tang 8.4% vs MLP(RV6). 384-dim return embedding 'nhan chim' 6-dim RV signal — hai khong gian feature orthogonal voi nhau."),
]

y = 1.25
for title, body in conclusions:
    rect(s, 0.3, y, 12.7, 1.15, fill=WHITE, line=ACCENT_BLUE, line_w=Pt(1))
    txt(s, title, 0.5, y + 0.04, 12.3, 0.4, size=12.5, bold=True, color=DARK_BLUE)
    txt(s, body,  0.5, y + 0.44, 12.3, 0.62, size=11, color=DARK_GRAY)
    y += 1.22


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Hướng cải thiện
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Huong Cai Thien Tiep Theo",
       "SOTA research 2024-2025: HAR-RV co the bi vuot voi dieu kien phu hop")
rect(s, 0, 1.1, 13.33, 6.4, fill=BG)

steps = [
    ("HAR-RS (Semivariance)", "Uu tien 1",
     "Thay the HAR-OLS bang semivariance decomposition: RV+ (return duong) va RV- (return am). "
     "Van la OLS, 5 tham so thay vi 3. Bat duoc leverage effect (downside vol du bao tot hon upside). "
     "Patton & Sheppard (2015): nhat quan vuot HAR o h=5,10,20 ngay — phu hop voi HORIZONS=[1,5,10,20] cua thi nghiem. ~10 dong code.",
     GREEN, "De, 1-2 gio"),
    ("GNN + QLIKE Loss", "Uu tien 2",
     "GNNHAR (Zhang et al., IJF 2025): GNN+QLIKE giam 13% MSE vs HAR tren US equity. "
     "Hien tai GNN dung MSE — doi sang QLIKE giup xu ly heteroskedasticity tot hon. "
     "Them cross-stock correlation graph tu daily returns. Khong can intraday data.",
     ACCENT_BLUE, "Vua, 1-2 ngay"),
    ("Foundation Model Fine-tuning", "Uu tien 3",
     "arxiv 2505.11163 (May 2025): TimesFM + incremental fine-tuning vuot HAR-RV thong ke (DM test) tren 21 chi so toan cau. "
     "Co the ap dung tuong tu cho Moirai2 voi LoRA adapter tren VN30 log-returns. "
     "Risk: phuc tap, co the ngoai pham vi luan van.",
     ACCENT_GOLD, "Dai han"),
]

y = 1.28
for title, priority, body, color, note in steps:
    rect(s, 0.3, y, 12.7, 1.72, fill=WHITE, line=color, line_w=Pt(2))
    rect(s, 0.3, y, 1.55, 1.72, fill=color)
    txt(s, priority, 0.3, y + 0.52, 1.55, 0.68,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 2.05, y + 0.05, 10.7, 0.42, size=13, bold=True, color=color)
    txt(s, body,  2.05, y + 0.5,  10.7, 1.1,  size=11, color=DARK_GRAY)
    txt(s, note,  10.8, y + 0.05, 2.1, 0.38,
        size=10, italic=True, color=RGBColor(0x77, 0x77, 0x77), align=PP_ALIGN.RIGHT)
    y += 1.88


# ═══════════════════════════════════════════════════════════════════════════════
OUT = "docs/slides_vn30_volatility_2026.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
